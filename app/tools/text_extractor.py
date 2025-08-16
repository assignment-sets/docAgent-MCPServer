# app/tools/text_extractor.py

import httpx
import io
import pandas as pd
import logging
from docx import Document
from pptx import Presentation
from app.schemas import FileExtractInput
from app.utils.Utils import Utils
import json
import os
from dotenv import load_dotenv
import tempfile
from langchain_community.document_loaders import PyMuPDFLoader
from typing import Callable

load_dotenv()

logger = logging.getLogger(__name__)


async def extract_text_from_pdf_file(input_data: FileExtractInput) -> str:
    """
    Extracts text from a PDF file (from a URL) using LangChain's PyMuPDF loader.
    """
    try:
        pdf_bytes = await Utils.download_file_as_stream(str(input_data.file_url))

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=True) as tmp_file:
            tmp_file.write(pdf_bytes)
            tmp_file.flush()

            loader = PyMuPDFLoader(tmp_file.name)
            docs = loader.load()

            # Step 4: Join all page_content strings
            full_text = "\n\n".join(
                doc.page_content for doc in docs if doc.page_content.strip()
            )

        if not full_text.strip():
            logger.warning("No text found in PDF: %s", input_data.file_url)
            raise ValueError("No text found in PDF.")

        logger.info("Successfully extracted text from PDF: %s", input_data.file_url)
        return full_text

    except Exception as e:
        logger.exception("Failed to extract text from PDF.")
        raise RuntimeError(f"PDF text extraction failed: {str(e)}")


async def extract_text_from_image_file(input_data: FileExtractInput) -> str:
    """
    Extracts text from an image URL using Google Vision API.
    The image must be publicly accessible (e.g., GCS signed URL or public web URL).
    """
    try:
        access_token = await Utils.get_google_access_token()

        headers = {
            "Authorization": f"Bearer {access_token}",
            "Content-Type": "application/json",
        }

        body = {
            "requests": [
                {
                    "image": {"source": {"imageUri": str(input_data.file_url)}},
                    "features": [{"type": "DOCUMENT_TEXT_DETECTION"}],
                }
            ]
        }

        async with httpx.AsyncClient(timeout=20) as client:
            response = await client.post(
                os.getenv("GOOGLE_VISION_IMAGE_OCR_API_URL"),
                headers=headers,
                data=json.dumps(body),
            )

        if response.status_code != 200:
            logger.error("Vision API returned error: %s", response.text)
            raise RuntimeError(f"Vision API error: {response.text}")

        result = response.json()
        extracted_text = (
            result["responses"][0].get("fullTextAnnotation", {}).get("text", "")
        )

        if not extracted_text.strip():
            logger.warning("No text detected in image: %s", input_data.file_url)
            raise ValueError("No text found in image.")

        logger.info("Successfully extracted text from image: %s", input_data.file_url)
        return extracted_text

    except Exception as e:
        logger.exception("Failed to extract text from image.")
        raise RuntimeError(f"Image text extraction failed: {str(e)}")


async def extract_text_from_pptx_file(input_data: FileExtractInput) -> str:
    """
    Extracts plain text from a PPTX file located at the given URL.
    Returns concatenated slide text as a single string.
    """
    try:
        file_stream = await Utils.download_file_as_stream(
            str(input_data.file_url), timeout=20
        )
        presentation = Presentation(io.BytesIO(file_stream))

        extracted_texts = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    extracted_texts.append(shape.text.strip())

        extracted_text = "\n".join(extracted_texts)

        if not extracted_text.strip():
            logger.warning("PPTX file is empty: %s", input_data.file_url)
            raise ValueError("PPTX file contains no extractable text.")

        logger.info("Successfully extracted text from PPTX: %s", input_data.file_url)
        return extracted_text

    except Exception as e:
        logger.exception("Failed to extract text from PPTX file.")
        raise RuntimeError(f"PPTX extraction failed: {str(e)}")


async def extract_text_from_docx_file(input_data: FileExtractInput) -> str:
    """
    Extracts plain text from a DOCX file located at the given URL.
    Returns the combined paragraph text as a string.
    """
    try:
        file_stream = await Utils.download_file_as_stream(str(input_data.file_url), 20)
        document = Document(io.BytesIO(file_stream))

        # Join all paragraphs with line breaks
        extracted_text = "\n".join(
            [para.text for para in document.paragraphs if para.text.strip()]
        )

        if not extracted_text.strip():
            logger.warning("DOCX file is empty: %s", input_data.file_url)
            raise ValueError("DOCX file contains no extractable text.")

        logger.info("Successfully extracted text from DOCX: %s", input_data.file_url)
        return extracted_text

    except Exception as e:
        logger.exception("Failed to extract text from DOCX file.")
        raise RuntimeError(f"DOCX extraction failed: {str(e)}")


async def extract_text_from_xlsx_file(input_data: FileExtractInput) -> str:
    """
    Extracts text from an XLSX file located at the given URL.
    Converts the first sheet into a JSON string suitable for LLM input.
    """
    try:
        file_stream = await Utils.download_file_as_stream(str(input_data.file_url), 20)
        excel_buffer = io.BytesIO(file_stream)
        df = pd.read_excel(excel_buffer, sheet_name=0)

        if df.empty:
            logger.warning("XLSX file is empty: %s", input_data.file_url)
            raise ValueError("XLSX file is empty or unreadable.")

        extracted_text = df.to_json(orient="records", force_ascii=False)

        logger.info("Successfully extracted JSON from XLSX: %s", input_data.file_url)
        return extracted_text

    except Exception as e:
        logger.exception("Failed to extract text from XLSX file.")
        raise RuntimeError(f"XLSX extraction failed: {str(e)}")


async def extract_text_from_csv_file(input_data: FileExtractInput) -> str:
    """
    Extracts text from a CSV file located at the given URL.
    Converts it into a JSON string format suitable for LLM input.
    """
    try:
        file_stream = await Utils.download_file_as_stream(str(input_data.file_url), 15)
        csv_buffer = io.StringIO(file_stream.decode("utf-8"))
        df = pd.read_csv(csv_buffer)

        if df.empty:
            logger.warning("CSV file is empty: %s", input_data.file_url)
            raise ValueError("CSV file is empty or unreadable.")

        extracted_text = df.to_json(orient="records", force_ascii=False)

        logger.info("Successfully extracted JSON from CSV: %s", input_data.file_url)
        return extracted_text

    except Exception as e:
        logger.exception("Failed to extract JSON from CSV file.")
        raise RuntimeError(f"CSV extraction failed: {str(e)}")


async def extract_text_from_plain_file(
    input_data: FileExtractInput,
) -> str:
    """
    Extracts raw text from .txt or .md file located at the given URL.
    Assumes UTF-8 encoding.
    """
    try:
        file_stream = await Utils.download_file_as_stream(str(input_data.file_url), 10)
        text_content = file_stream.decode("utf-8").strip()

        if not text_content:
            logger.warning("Empty content extracted from file: %s", input_data.file_url)
            raise ValueError("File is empty or unreadable.")

        logger.info("Successfully extracted text from: %s", input_data.file_url)
        return text_content

    except Exception as e:
        logger.exception("Unexpected error during plain text extraction.")
        raise RuntimeError(f"Text extraction failed: {str(e)}")


EXTENSION_HANDLER_MAP: dict[str, Callable[[FileExtractInput], str]] = {
    "pdf": extract_text_from_pdf_file,
    "png": extract_text_from_image_file,
    "jpg": extract_text_from_image_file,
    "jpeg": extract_text_from_image_file,
    "webp": extract_text_from_image_file,
    "docx": extract_text_from_docx_file,
    "csv": extract_text_from_csv_file,
    "xlsx": extract_text_from_xlsx_file,
    "pptx": extract_text_from_pptx_file,
    "txt": extract_text_from_plain_file,
    "md": extract_text_from_plain_file,
}


async def extract_text(input_data: FileExtractInput) -> str:
    """
    Uses provided file_type to select the correct extractor
    """
    try:
        ext = input_data.file_type.lower().lstrip(".")  # clean 'pdf' or '.pdf'

        if ext in EXTENSION_HANDLER_MAP:
            extractor_func = EXTENSION_HANDLER_MAP[ext]
            logger.info(f"Dispatching to extractor for .{ext} file.")
            return await extractor_func(input_data)
        else:
            logger.warning(f"Unsupported file type: .{ext}")
            raise ValueError(f"Unsupported file type: .{ext}")

    except Exception as e:
        logger.exception("Text extraction failed.")
        raise RuntimeError(f"Failed to extract text: {str(e)}")


if __name__ == "__main__":
    import asyncio

    async def local_test():
        res = await extract_text(
            FileExtractInput(
                file_url="https://storage.googleapis.com/doc-agent-buck-1/temp/Iris.csv"
            )
        )
        print(res)

    asyncio.run(local_test())
